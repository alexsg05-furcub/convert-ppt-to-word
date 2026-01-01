from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.responses import FileResponse
from fastapi.middleware.cors import CORSMiddleware
from pptx import Presentation
from docx import Document
import uuid
import os

app = FastAPI()

# Permitir CORS para tu HTML
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Puedes cambiar a ["http://localhost:5500"] si quieres restringir
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

UPLOAD_DIR = "uploads"
OUTPUT_DIR = "output"

os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(OUTPUT_DIR, exist_ok=True)


def clean_text(text: str) -> str:
    """
    Limpia saltos de línea y espacios extra.
    """
    return ' '.join(text.split()).strip()


def pptx_to_docx_clean(pptx_path: str, docx_path: str):
    """
    Convierte un PPTX a DOCX limpio:
    - Cada diapositiva con un título "Diapositiva X"
    - Párrafos limpios
    - Mantiene estructura básica de viñetas
    """
    presentation = Presentation(pptx_path)
    document = Document()

    for slide_index, slide in enumerate(presentation.slides, start=1):
        document.add_heading(f"Diapositiva {slide_index}", level=1)

        for shape in slide.shapes:
            if shape.has_text_frame:
                # Itera párrafo por párrafo
                for para in shape.text_frame.paragraphs:
                    text = clean_text(para.text)
                    if text:  # Ignora párrafos vacíos
                        document.add_paragraph(text)

        # Línea separadora entre diapositivas
        document.add_paragraph("\n" + "-" * 40 + "\n")

    document.save(docx_path)


@app.post("/convert")
async def convert_ppt_to_word(file: UploadFile = File(...)):
    if not file.filename.endswith((".pptx", ".ppt")):
        raise HTTPException(status_code=400, detail="Archivo no válido")

    file_id = str(uuid.uuid4())
    pptx_path = os.path.join(UPLOAD_DIR, f"{file_id}.pptx")
    docx_path = os.path.join(OUTPUT_DIR, f"{file_id}.docx")

    with open(pptx_path, "wb") as f:
        f.write(await file.read())

    try:
        pptx_to_docx_clean(pptx_path, docx_path)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error al procesar el archivo: {e}")

    return FileResponse(
        docx_path,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        filename="resultado.docx"
    )
